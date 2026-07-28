"""Verificacion Fase 6: cifras de las slides == artefactos del modelo final."""
import json
import re
from pptx import Presentation

BASE = r"C:\Users\bryan\ia\proyecto_iaaa\Proyecto\ProyectoIAAA"
mj = json.load(open(BASE + r"\Data\ml_out\lstm_fantasma_final\metrics_test.json", encoding="utf-8"))
bj = json.load(open(BASE + r"\Data\ml_out\lstm_fantasma_final\binary_metrics_test.json", encoding="utf-8"))

exp = []
for t in ("y3", "y5", "y10", "y15"):
    r = mj["regression"][t]
    b = bj[t]
    exp += [
        f"{r['mae']:.2f}", f"{r['rmse']:.2f}",
        f"{b['accuracy']:.2f}", f"{b['precision']:.2f}", f"{b['recall']:.2f}", f"{b['f1']:.2f}",
    ]
exp += ["1.49", "1.84", "0.87"]

p = Presentation(BASE + r"\docs\material_profesor\PRESENTACION_FINAL_ML.pptx")
alltxt = " ".join(sh.text_frame.text for s in p.slides for sh in s.shapes if sh.has_text_frame)
missing = [v for v in dict.fromkeys(exp) if v not in alltxt]
print("valores esperados (2dec):", len(set(exp)))
print("AUSENTES en pptx:", missing if missing else "ninguno")

conf = []
for t in ("y3", "y5", "y10", "y15"):
    cj = bj[t]["confusion"]
    conf += [cj["tp"], cj["fp"], cj["tn"], cj["fn"]]
flat = re.sub(r"[^0-9]", " ", alltxt)
flat = " ".join(flat.split()).split()
miss_c = [str(v) for v in conf if str(v) not in flat]
print("confusion ausente:", miss_c if miss_c else "ninguna")

stale = [v for v in ["0.91", "1.31", "2.81", "3.58", "0.48", "0.57", "0.44",
                     "1 051", "1 402", "1 547", "32 unidades", "20 épocas",
                     "lstm_fantasma/fantasma_lstm.params", "train_fantasma_lstm.pl"] if v in alltxt]
print("rastro v1 en pptx:", stale if stale else "ninguno")

html = open(BASE + r"\docs\material_profesor\PRESENTACION_FINAL_ML.html", encoding="utf-8").read()
missing_h = [v for v in dict.fromkeys(exp) if v not in html]
print("AUSENTES en html:", missing_h if missing_h else "ninguno")
stale_h = [v for v in ["0.91", "1.31", "2.81", "3.58", "0.48", "0.57", "0.44",
                       "32 unidades", "20 épocas"] if v in html]
print("rastro v1 en html:", stale_h if stale_h else "ninguno")

for doc in ("GUION_PRESENTACION_FINAL.md", "CHECKLIST_DEMO_FINAL.md"):
    txt = open(BASE + rf"\docs\material_profesor\{doc}", encoding="utf-8").read()
    v1 = [v for v in ["lstm_fantasma/fantasma_lstm.params", "lstm_fantasma/metrics_test.json",
                      "lstm_fantasma/preds_test.csv", "lstm_fantasma/binary_metrics_test.json",
                      "train_fantasma_lstm.pl --eval-only", "0.91", "1 051"] if v in txt]
    print(f"rastro v1 en {doc}:", v1 if v1 else "ninguno")
