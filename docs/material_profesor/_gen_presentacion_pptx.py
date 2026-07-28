"""Generate PRESENTACION_FINAL_ML.pptx — editable monochrome deck, title cover, no headers/footers."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUT = r"C:\Users\bryan\ia\proyecto_iaaa\Proyecto\ProyectoIAAA\docs\material_profesor\PRESENTACION_FINAL_ML.pptx"

INK = RGBColor(0x16, 0x16, 0x16)
INK2 = RGBColor(0x3D, 0x3D, 0x3D)
INK3 = RGBColor(0x6B, 0x6B, 0x6B)
PAPER = RGBColor(0xFB, 0xFB, 0xF9)
HAIR = RGBColor(0xD9, 0xD9, 0xD4)
WHITE = RGBColor(0xFB, 0xFB, 0xF9)
DIM = RGBColor(0x9A, 0x9A, 0x90)

SERIF = "Georgia"
SANS = "Segoe UI"
MONO = "Consolas"

EMU_IN = 914400
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = 13.333, 7.5
blank = prs.slide_layouts[6]


def rect(slide, x, y, w, h, color):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def line(slide, x, y, w, weight=2.0, color=INK):
    ln = slide.shapes.add_connector(1, Inches(x), Inches(y), Inches(x + w), Inches(y))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    ln.shadow.inherit = False
    return ln


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15, space_after=0):
    """runs: list of paragraphs; each paragraph = list of (txt, font, size, bold, color)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.line_spacing = line_spacing
        if space_after:
            p.space_after = Pt(space_after)
        for (t, font, size, bold, color) in para:
            r = p.add_run()
            r.text = t
            r.font.name = font
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
    return tb


def P(t, font=SANS, size=12, bold=False, color=INK2):
    return [(t, font, size, bold, color)]


def pageno(slide, n, dark=False):
    text(slide, SW - 1.0, SH - 0.52, 0.6, 0.3,
         [P(f"{n:02d}", MONO, 10, False, DIM if not dark else RGBColor(0x8A, 0x8A, 0x80))],
         align=PP_ALIGN.RIGHT)


def bg(slide, color=PAPER):
    rect(slide, -0.1, -0.1, SW + 0.2, SH + 0.2, color)


def kicker_rule(slide, kicker, n):
    text(slide, 0.6, 0.42, 10, 0.3, [P(kicker.upper(), SANS, 9, True, INK3)])
    # small tracking via spaces not available; keep simple
    line(slide, 0.6, 0.78, SW - 1.2, 2.4, INK)
    pageno(slide, n)


def table(slide, x, y, w, headers, rows, caption, col_fracs=None, row_h=0.52, head_h=0.48,
          base_size=13, dark=False):
    ink = WHITE if dark else INK
    ink2 = RGBColor(0xC9, 0xC9, 0xC0) if dark else INK2
    hair = RGBColor(0x3D, 0x3D, 0x3D) if dark else HAIR
    ink3 = RGBColor(0x9A, 0x9A, 0x90) if dark else INK3
    ncols = len(headers)
    if col_fracs is None:
        col_fracs = [0.36] + [(0.64) / (ncols - 1)] * (ncols - 1)
    colw = [w * f for f in col_fracs]
    text(slide, x, y, w, 0.26, [P(caption.upper(), SANS, 8.5, True, ink3)])
    top = y + 0.34
    line(slide, x, top, w, 2, ink)
    hy = top + head_h
    xx = x
    for i, hname in enumerate(headers):
        al = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.RIGHT
        text(slide, xx + 0.06, top + 0.1, colw[i] - 0.12, head_h - 0.12,
             [P(hname, SANS, 9.5, True, ink)], align=al)
        xx += colw[i]
    line(slide, x, hy, w, 2, ink)
    yy = hy
    for r_i, row in enumerate(rows):
        yy = hy + r_i * row_h
        xx = x
        for i, val in enumerate(row):
            al = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.RIGHT
            font = SANS if i == 0 else MONO
            color = ink if i == 0 else ink2
            bold = (i == 0)
            text(slide, xx + 0.06, yy + 0.1, colw[i] - 0.12, row_h - 0.12,
                 [P(val, font, base_size, bold, color)], align=al)
            xx += colw[i]
        line(slide, x, yy + row_h, w, 0.9, hair)
    line(slide, x, yy + row_h, w, 2, ink)


# ============ 1 · PORTADA (título grande, sin kicker ni footer) ============
s = prs.slides.add_slide(blank)
bg(s)
line(s, 0.6, 0.9, SW - 1.2, 3, INK)
text(s, 0.6, 1.35, SW - 1.2, 2.6, [
    [("Inteligencia Artificial", SERIF, 44, True, INK)],
    [("y Aprendizaje Automático", SERIF, 44, True, INK)],
    [("Proyecto Final", SERIF, 54, True, INK)],
], line_spacing=1.05)
text(s, 0.6, 4.25, 11.5, 0.9, [
    [("Predicción del comportamiento del fantasma con un modelo LSTM, ", SANS, 14, False, INK2),
     ("a partir de la estructura y la liquidez del motor de charting.", SANS, 14, False, INK2)],
], line_spacing=1.35)
stats = [
    ("7 649", "muestras de entrenamiento (abril–junio)"),
    ("2 391", "muestras de prueba (1–24 de julio)"),
    ("86", "variables de entrada normalizadas"),
    ("4", "salidas: y3 · y5 · y10 · y15"),
]
for i, (v, k) in enumerate(stats):
    x = 0.6 + i * 3.1
    line(s, x, 5.55, 2.5, 2.4, INK)
    text(s, x, 5.7, 2.7, 0.7, [P(v, SERIF, 34, True, INK)])
    text(s, x, 6.45, 2.7, 0.6, [P(k, SANS, 10.5, False, INK2)], line_spacing=1.2)
pageno(s, 1)

# ============ 2 · OBJETIVO ============
s = prs.slides.add_slide(blank)
bg(s)
kicker_rule(s, "01 · Objetivo del modelo", 2)
text(s, 0.6, 1.0, 12, 0.8, [P("Cuatro salidas, una por horizonte", SERIF, 30, True, INK)])
text(s, 0.6, 1.75, 7.1, 1.3, [
    P("Cada vez que el fantasma aparece o se reubica se toma una muestra. El modelo predice "
      "cuántos rastros dejará hacia adelante en cuatro horizontes distintos.", SANS, 13, False, INK2),
], line_spacing=1.3)
# deflist left
line(s, 0.6, 3.0, 6.9, 1.8, INK)
text(s, 0.6, 3.12, 6, 0.3, [P("DEFINICIÓN DE LAS SALIDAS", SANS, 8.5, True, INK3)])
defs = [
    ("y3", "rastros que deja el fantasma en los próximos 3 minutos"),
    ("y5", "rastros en los próximos 5 minutos"),
    ("y10", "rastros en los próximos 10 minutos"),
    ("y15", "rastros en los próximos 15 minutos"),
]
for i, (lab, txt) in enumerate(defs):
    yy = 3.5 + i * 0.78
    text(s, 0.6, yy, 1.0, 0.5, [P(lab, MONO, 14, True, INK)])
    text(s, 1.7, yy, 5.8, 0.7, [P(txt, SANS, 13, False, INK2)], line_spacing=1.15)
    line(s, 0.6, yy + 0.62, 6.9, 0.9, HAIR)
# right box
line(s, 8.2, 3.0, 4.5, 1.8, INK)
text(s, 8.2, 3.12, 4.3, 0.3, [P("ENTRADA DEL MODELO", SANS, 8.5, True, INK3)])
text(s, 8.2, 3.5, 4.5, 3.2, [
    [("Una ventana de ", SANS, 12.5, False, INK2), ("5 muestras consecutivas", SANS, 12.5, True, INK),
     (", cada una con ", SANS, 12.5, False, INK2), ("86 variables", SANS, 12.5, True, INK),
     (" calculadas en la vela siguiente al evento:", SANS, 12.5, False, INK2)],
    [("distancias en PIPs a niveles de liquidez y estructura (order blocks, FVG, Fibonacci, "
      "AVWAP, perfil de volumen, rupturas de estructura), en tres temporalidades: ",
      SANS, 12.5, False, INK2), ("1m, 10m y 1h.", SANS, 12.5, True, INK)],
], line_spacing=1.3, space_after=10)

# ============ 3 · MUESTREO ============
s = prs.slides.add_slide(blank)
bg(s)
kicker_rule(s, "02 · Construcción de las muestras", 3)
text(s, 0.6, 1.0, 12, 0.8, [P("El disparo es el evento del fantasma", SERIF, 30, True, INK)])
steps = [
    ("01", "Detección del evento",
     "Se recorre el histórico 1m en modo Replay y se registra cada aparición o reubicación del fantasma provisional."),
    ("02", "Contexto causal",
     "Las variables se calculan mirando solo hacia atrás, con la información disponible hasta la vela "
     "siguiente al evento. Nada mira el futuro."),
    ("03", "Etiquetado automático",
     "Las salidas y3/y5/y10/y15 se obtienen contando los rastros que efectivamente aparecen después del evento."),
]
for i, (num, head, body) in enumerate(steps):
    yy = 2.1 + i * 1.35
    text(s, 0.6, yy, 0.6, 0.4, [P(num, MONO, 11, True, INK3)])
    text(s, 1.2, yy, 6.3, 0.35, [P(head, SANS, 14, True, INK)])
    text(s, 1.2, yy + 0.36, 6.3, 0.85, [P(body, SANS, 12, False, INK2)], line_spacing=1.25)
# right volume
line(s, 8.2, 2.0, 4.5, 1.8, INK)
text(s, 8.2, 2.12, 4.3, 0.3, [P("VOLUMEN DEL DATASET", SANS, 8.5, True, INK3)])
vol = [
    ("Train", "7 649 muestras · abr–jun (88 736 velas)"),
    ("Test", "2 391 muestras · 1–24 jul (24 179 velas)"),
    ("Secuencias", "7 645 train / 2 387 test (ventana 5)"),
]
for i, (lab, txt) in enumerate(vol):
    yy = 2.5 + i * 0.72
    text(s, 8.2, yy, 1.3, 0.5, [P(lab, MONO, 11.5, True, INK)])
    text(s, 9.5, yy, 3.2, 0.65, [P(txt, SANS, 11.5, False, INK2)], line_spacing=1.15)
    line(s, 8.2, yy + 0.58, 4.5, 0.9, HAIR)
text(s, 8.2, 4.85, 4.5, 0.9, [
    P("El comportamiento del fantasma replica el indicador implementado en el motor de charting.",
      SANS, 10.5, False, INK3)], line_spacing=1.25)

# ============ 4 · PROCESO ============
s = prs.slides.add_slide(blank)
bg(s)
kicker_rule(s, "03 · Proceso", 4)
text(s, 0.6, 1.0, 12, 0.8, [P("Del CSV al modelo", SERIF, 30, True, INK)])
cards = [
    ("1 · EXTRACCIÓN",
     "Script sin interfaz gráfica recorre el CSV en Replay y, en cada evento del fantasma, "
     "guarda una fila con las 86 variables y las 4 salidas futuras.",
     "7 649 filas train · 2 391 test"),
    ("2 · NORMALIZACIÓN",
     "Estandarización z-score con media y desviación calculadas solo con train; "
     "el test se transforma con los mismos parámetros.",
     "Tiempo y metadatos no entrenan"),
    ("3 · ENTRENAMIENTO",
     "LSTM de una capa (48 unidades, dropout 0.2), ventana de 5 muestras, 4 salidas, "
     "error cuadrático medio, optimizador Adam y parada temprana con validación.",
     "Pesos guardados en disco"),
]
for i, (head, body, note) in enumerate(cards):
    x = 0.6 + i * 4.2
    line(s, x, 2.1, 3.7, 1.8, INK)
    text(s, x, 2.22, 3.6, 0.3, [P(head, SANS, 8.5, True, INK3)])
    text(s, x, 2.6, 3.7, 2.4, [P(body, SANS, 12, False, INK2)], line_spacing=1.3)
    text(s, x, 5.5, 3.7, 0.4, [P(note, SANS, 10.5, False, INK3)])
text(s, 0.6, 6.3, 12, 0.5, [
    P("Grid de 8 configuraciones (~8 min), elegida por validación · la demostración carga los pesos, no reentrena",
      SANS, 11, False, INK3)])

# ============ 5 · REGRESIÓN ============
s = prs.slides.add_slide(blank)
bg(s)
kicker_rule(s, "04 · Resultados sobre julio (holdout)", 5)
text(s, 0.6, 1.0, 12, 0.8, [P("Error en el conteo de rastros", SERIF, 30, True, INK)])
table(s, 0.6, 1.9, 6.6,
      ["Salida", "MAE", "RMSE"],
      [("y3 · 3 min", "0.82", "1.01"),
       ("y5 · 5 min", "1.15", "1.41"),
       ("y10 · 10 min", "1.75", "2.15"),
       ("y15 · 15 min", "2.24", "2.77")],
      "Regresión · n = 2 387 secuencias de test", row_h=0.6, base_size=14)
text(s, 8.0, 2.3, 4.7, 2.4, [
    [("En la ventana corta el modelo se equivoca en ", SANS, 13, False, INK2),
     ("menos de un rastro", SANS, 13, True, INK), (" en promedio (MAE 0.82 en y3).", SANS, 13, False, INK2)],
    [("El error crece con el horizonte, de forma esperada: a 15 minutos hay más rastros posibles "
      "y más incertidumbre.", SANS, 13, False, INK2)],
    [("El error medio de las cuatro ventanas es ", SANS, 13, False, INK2),
     ("1.49 rastros", SANS, 13, True, INK),
     (", un 17% menos que la primera versión entrenada del modelo.", SANS, 13, False, INK2)],
], line_spacing=1.35, space_after=10)
text(s, 0.6, 5.35, 7, 0.5, [
    P("MAE: error medio absoluto en rastros · RMSE penaliza errores grandes", SANS, 10.5, False, INK3)])

# ============ 6 · BINARIA ============
s = prs.slides.add_slide(blank)
bg(s)
kicker_rule(s, "04 · Vista binaria", 6)
text(s, 0.6, 1.0, 12, 0.8, [P("¿Aparece al menos un rastro?", SERIF, 30, True, INK)])
table(s, 0.6, 1.9, 7.6,
      ["Salida", "Accuracy", "Precision", "Recall", "F1"],
      [("y3", "0.71", "0.69", "0.98", "0.81"),
       ("y5", "0.75", "0.75", "0.99", "0.85"),
       ("y10", "0.81", "0.81", "0.99", "0.89"),
       ("y15", "0.83", "0.84", "0.99", "0.91")],
      "Positivo = ≥1 rastro en la ventana · n = 2 387", row_h=0.6, base_size=14)
text(s, 8.8, 2.3, 3.9, 2.8, [
    [("El modelo casi no se pierde una ventana con actividad: detecta ", SANS, 12.5, False, INK2),
     ("98–99%", SANS, 12.5, True, INK),
     (" de los casos con al menos un rastro (recall).", SANS, 12.5, False, INK2)],
    [("Cuando anuncia actividad a 10–15 minutos acierta 81–84% de las veces (precision); "
      "el F1 promedio de las cuatro ventanas es ", SANS, 12.5, False, INK2),
     ("0.87", SANS, 12.5, True, INK), (".", SANS, 12.5, False, INK2)],
], line_spacing=1.3, space_after=10)
text(s, 0.6, 5.35, 7, 0.5, [
    P("La salida continua se binariza con umbral 0.5", SANS, 10.5, False, INK3)])

# ============ 7 · CONFUSIÓN ============
s = prs.slides.add_slide(blank)
bg(s)
kicker_rule(s, "04 · Matriz de confusión", 7)
text(s, 0.6, 1.0, 12, 0.8, [P("Detalle por ventana", SERIF, 30, True, INK)])
table(s, 0.6, 1.9, 7.6,
      ["Salida", "TP", "FP", "TN", "FN"],
      [("y3", "1 513", "668", "174", "32"),
       ("y5", "1 675", "563", "126", "23"),
       ("y10", "1 862", "427", "81", "17"),
       ("y15", "1 939", "380", "46", "22")],
      "Positivo = ≥1 rastro · n = 2 387 por fila", row_h=0.6, base_size=14)
line(s, 8.8, 2.1, 3.9, 1.8, INK)
text(s, 8.8, 2.22, 3.7, 0.3, [P("LECTURA", SANS, 8.5, True, INK3)])
conf_defs = [
    ("TP", "predijo rastro y sí hubo"),
    ("FP", "predijo rastro y no hubo"),
    ("TN", "predijo que no y no hubo"),
    ("FN", "predijo que no y sí hubo"),
]
for i, (lab, txt) in enumerate(conf_defs):
    yy = 2.6 + i * 0.55
    text(s, 8.8, yy, 0.7, 0.4, [P(lab, MONO, 12.5, True, INK)])
    text(s, 9.5, yy, 3.2, 0.4, [P(txt, SANS, 12, False, INK2)])
text(s, 8.8, 5.0, 3.9, 1.0, [
    P("Casi no deja pasar actividad real: solo 17–32 falsos negativos por ventana, "
      "de 1 545–1 961 casos con rastro.",
      SANS, 10.5, False, INK3)], line_spacing=1.25)

# ============ 8 · DEMO ============
s = prs.slides.add_slide(blank)
bg(s)
kicker_rule(s, "05 · Demostración en vivo", 8)
text(s, 0.6, 1.0, 12, 0.8, [P("Cargar el modelo y comparar", SERIF, 30, True, INK)])
text(s, 0.6, 1.8, 12.1, 1.0, [
    [("El programa carga los pesos entrenados (", SANS, 13, False, INK2),
     ("no reentrena", SANS, 13, True, INK),
     ("), arma las ventanas de 5 muestras sobre el test ya normalizado e imprime el conteo real "
      "contra la predicción para varios puntos de julio.", SANS, 13, False, INK2)],
], line_spacing=1.3)
cmd = rect(s, 0.6, 2.9, 12.13, 1.5, INK)
text(s, 1.0, 3.12, 11.5, 1.2, [
    [("# entorno: WSL Fedora · raíz del proyecto", MONO, 10.5, False, RGBColor(0x8A, 0x8A, 0x80))],
    [("cd /mnt/c/Users/bryan/ia/proyecto_iaaa/Proyecto/ProyectoIAAA", MONO, 12.5, True, WHITE)],
    [("perl -I. scripts/demo_fantasma_predict.pl --n 8", MONO, 12.5, True, WHITE)],
], line_spacing=1.45)
demo_steps = [
    ("PASO 1", "Carga el test normalizado y los pesos .params"),
    ("PASO 2", "Construye ventanas de 5 muestras consecutivas"),
    ("PASO 3", "Imprime TRUE vs PRED para y3 · y5 · y10 · y15"),
]
for i, (head, body) in enumerate(demo_steps):
    x = 0.6 + i * 4.2
    line(s, x, 4.9, 3.7, 1.8, INK)
    text(s, x, 5.02, 3.6, 0.3, [P(head, SANS, 8.5, True, INK3)])
    text(s, x, 5.36, 3.7, 1.0, [P(body, SANS, 12, False, INK2)], line_spacing=1.25)
text(s, 0.6, 6.7, 12, 0.4, [
    P("Ejecución en pocos segundos · respaldo: métricas y predicciones guardadas en disco",
      SANS, 10.5, False, INK3)])

# ============ 9 · CIERRE (oscuro) ============
s = prs.slides.add_slide(blank)
bg(s, INK)
line(s, 0.6, 0.78, SW - 1.2, 2.4, WHITE)
text(s, 0.6, 1.0, 12, 0.8, [P("Qué queda demostrado", SERIF, 30, True, WHITE)])
close_items = [
    ("01", "El disparo es correcto",
     "Cada muestra corresponde a un evento real del fantasma, con contexto estrictamente causal."),
    ("02", "El pipeline es reproducible",
     "Extracción, normalización, entrenamiento y evaluación quedan como scripts ejecutables."),
    ("03", "El modelo predice en vivo",
     "Carga los pesos y compara contra datos que nunca vio en entrenamiento."),
]
for i, (num, head, body) in enumerate(close_items):
    yy = 2.1 + i * 1.3
    text(s, 0.6, yy, 0.6, 0.4, [P(num, MONO, 11, True, DIM)])
    text(s, 1.2, yy, 6.3, 0.35, [P(head, SANS, 14, True, WHITE)])
    text(s, 1.2, yy + 0.36, 6.3, 0.85,
         [P(body, SANS, 12, False, RGBColor(0xC9, 0xC9, 0xC0))], line_spacing=1.25)
table(s, 8.2, 2.0, 4.5,
      ["Pieza", "Valor"],
      [("Dataset train/test", "7 649 / 2 391"),
       ("Variables", "86"),
       ("MAE prom · 4 ventanas", "1.49"),
       ("F1 prom · binario", "0.87")],
      "Resumen", row_h=0.56, base_size=12.5, dark=True, col_fracs=[0.58, 0.42])
pageno(s, 9, dark=True)

prs.save(OUT)
print("Wrote", OUT)
